import sys, io, os, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Emu
from lxml import etree

DOCS  = r"C:\Users\Sergej\Documents\GitHub\tkdBot\MyHouse\docs"
PPTX  = os.path.join(DOCS, "Я устал, я не хочу работать бесплатно (2).pptx")
OUT   = os.path.join(DOCS, "Я устал, я не хочу работать бесплатно (2).pptx")

KPI_PNG   = os.path.join(DOCS, "slide6_new_kpi.png")
CHART_PNG = os.path.join(DOCS, "slide6_new_chart.png")

prs = Presentation(PPTX)

# ── SLIDE 6: replace embedded images ──────────────────────────────────────────
slide6 = prs.slides[5]

pic_shapes = [s for s in slide6.shapes if s.shape_type == 13]
assert len(pic_shapes) == 2, f"Expected 2 pictures, got {len(pic_shapes)}"

replacements = [KPI_PNG, CHART_PNG]

for shape, new_img_path in zip(pic_shapes, replacements):
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    # Add new picture at same geometry
    new_shape = slide6.shapes.add_picture(new_img_path, left, top, width, height)
    # Move new shape's XML element to same position in z-order as old shape
    old_el = shape._element
    new_el = new_shape._element
    parent = old_el.getparent()
    idx = list(parent).index(old_el)
    parent.remove(old_el)
    parent.remove(new_el)
    parent.insert(idx, new_el)
    print(f"Replaced {shape.name} with {os.path.basename(new_img_path)}")


# ── SLIDE 8: text replacement ──────────────────────────────────────────────────
slide8 = prs.slides[7]

REPLACEMENTS = {
    "30-50 объектов":   "15–25 крупных контрактов",
    "OEM лицензия":     "Партнёрский канал (КРОК, Rubytech)",
    "978k":             "Pre-seed: 5–15M",
    "финансирование (₽)": "цель Pre-seed раунда",
}

def replace_text_in_shape(shape, replacements):
    if hasattr(shape, "shapes"):
        for s in shape.shapes:
            replace_text_in_shape(s, replacements)
        return
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        full_text = para.text
        for old, new in replacements.items():
            if old in full_text:
                # Put replacement in first run, clear others
                if para.runs:
                    para.runs[0].text = full_text.replace(old, new)
                    for run in para.runs[1:]:
                        run.text = ""
                    print(f"  Replaced: '{old}' → '{new}'  (in: '{full_text.strip()}')")
                break

print("\nSlide 8 text replacements:")
for shape in slide8.shapes:
    replace_text_in_shape(shape, REPLACEMENTS)

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"\nSaved: {OUT}")
